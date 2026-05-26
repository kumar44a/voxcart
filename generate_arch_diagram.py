#!/usr/bin/env python3
"""
VoxCart — Architecture Diagram Generator
Produces a slide-ready PNG and embeds it in VoxCart_Architecture_Slide.pptx
"""

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch
import os

OUT_PNG  = os.path.join(os.path.dirname(__file__), 'assets', 'architecture_diagram.png')
OUT_PPTX = os.path.join(os.path.dirname(__file__), 'VoxCart_Architecture_Slide.pptx')
os.makedirs(os.path.dirname(OUT_PNG), exist_ok=True)

# ── Figure (13.33 × 7.5 → same ratio as 16:9 slide) ─────────────────────────
fig, ax = plt.subplots(figsize=(13.33, 7.5))
ax.set_xlim(0, 13.33)
ax.set_ylim(0, 7.5)
ax.axis('off')
fig.patch.set_facecolor('#F4F7FC')
plt.subplots_adjust(left=0, right=1, top=1, bottom=0)

# ── Helpers ───────────────────────────────────────────────────────────────────
def rbox(x, y, w, h, fc, ec, lw=1.5, z=2):
    ax.add_patch(FancyBboxPatch(
        (x, y), w, h, boxstyle='round,pad=0.04',
        facecolor=fc, edgecolor=ec, linewidth=lw, zorder=z, clip_on=False))

def t(x, y, s, sz=9, bold=False, col='#111133', ha='center', va='center', z=5):
    ax.text(x, y, s, fontsize=sz, fontweight='bold' if bold else 'normal',
            color=col, ha=ha, va=va, zorder=z)

def arr(x1, y1, x2, y2, col='#1A3399', lw=1.8, rad=0.0):
    ax.annotate('', xy=(x2, y2), xytext=(x1, y1),
                arrowprops=dict(arrowstyle='->', color=col, lw=lw,
                                connectionstyle=f'arc3,rad={rad}',
                                mutation_scale=14),
                zorder=8)

def badge(x, y, s, sz=7.5, col='#1A3399'):
    ax.text(x, y, s, fontsize=sz, fontweight='bold', color=col,
            ha='center', va='center', zorder=9,
            bbox=dict(boxstyle='round,pad=0.28', fc='#FFFFFF',
                      ec='#99AACC', alpha=0.95, lw=0.9))

# ═══════════════════════════════════════════════════════════════════════════════
# TIER 1 — USER'S BROWSER     y 5.75 – 7.33
# ═══════════════════════════════════════════════════════════════════════════════
rbox(0.25, 5.75, 12.85, 1.58, '#D3E9F7', '#2472A4', lw=2.2, z=1)
t(0.60, 7.18, "USER'S BROWSER", 9.5, True, '#1A4A7A', ha='left')

browser_labels = [
    'Microphone\nCapture',
    'Avatar + State\nVisualizer',
    'Transcript\nPanel',
    'Audio Output\nPlayback',
]
for i, lbl in enumerate(browser_labels):
    bx = 0.44 + i * 3.18
    rbox(bx, 5.88, 2.96, 1.30, '#FFFFFF', '#85B5D5', lw=1.0, z=3)
    t(bx + 1.48, 6.54, lbl, 8.2, False, '#1A3A6A')

# ═══════════════════════════════════════════════════════════════════════════════
# TIER 2-L — api.py (Flask)    x 0.25–3.75,  y 3.55–5.45
# ═══════════════════════════════════════════════════════════════════════════════
rbox(0.25, 3.55, 3.50, 1.90, '#FFF2DC', '#D48B10', lw=2.0)
t(2.00, 4.78, 'api.py',           11.5, True,  '#8B4500')
t(2.00, 4.46, 'Flask  :5001',      9.0, False, '#7A5030')
t(2.00, 4.17, 'JWT token minting', 8.2, False, '#AA8866')

# ═══════════════════════════════════════════════════════════════════════════════
# TIER 2-R — LiveKit Cloud     x 4.60–10.20,  y 3.35–5.45
# ═══════════════════════════════════════════════════════════════════════════════
rbox(4.60, 3.35, 5.60, 2.10, '#CDDEFF', '#1A44C8', lw=2.5)
t(7.40, 4.78, 'LIVEKIT CLOUD',               12.5, True,  '#0B2EA8')
t(7.40, 4.44, 'WebRTC Media Server',           9.5, False, '#334488')
t(7.40, 4.14, 'Room  ·  Tracks  ·  Dispatch',  8.5, False, '#556699')

# ═══════════════════════════════════════════════════════════════════════════════
# TIER 3 — agent.py (VoxCart Agent)     y 0.15 – 3.10
# ═══════════════════════════════════════════════════════════════════════════════
rbox(0.25, 0.15, 12.85, 2.95, '#FAF4E8', '#C07820', lw=2.2, z=1)
t(0.60, 2.90, 'agent.py  (VoxCart Agent)', 9.5, True, '#885500', ha='left')

# Silero VAD
rbox(0.44, 1.83, 1.90, 0.93, '#E4F6E4', '#28A745', lw=1.5, z=3)
t(1.39, 2.37, 'Silero',  10.0, True,  '#1B6B3A')
t(1.39, 2.07, 'VAD',      8.5, False, '#2D8C50')

# OpenAI Whisper STT
rbox(2.57, 1.83, 2.25, 0.93, '#E4ECFF', '#3060CC', lw=1.5, z=3)
t(3.70, 2.37, 'OpenAI Whisper',  9.5, True,  '#1A3B8C')
t(3.70, 2.07, 'STT',             8.5, False, '#3355AA')

# GPT-4.1-mini LLM block
rbox(5.04, 1.27, 7.68, 1.49, '#F0E8FF', '#7020C0', lw=1.8, z=3)
t(8.88, 2.57, 'GPT-4.1-mini  (LLM)',                      10.5, True,  '#4A0080')
t(8.88, 2.29, 'E-Commerce System Prompt  +  RAG Context',   8.2, False, '#665588')

# Function-calling tools sub-box (inside LLM block)
rbox(5.19, 1.28, 7.36, 0.78, '#ECFBEC', '#40A840', lw=1.0, z=4)
t(8.87, 1.84, 'Function-Calling Tools', 8.0, True, '#1A5C1A')
t(8.87, 1.57,
  'get_order_status   ·   lookup_product   ·   get_returns_policy   ·   search_faq   ·   get_current_datetime',
  7.2, False, '#1A5C1A')

# Cartesia TTS
rbox(5.04, 0.27, 4.25, 0.73, '#FFECEC', '#C02020', lw=1.8, z=3)
t(7.17, 0.64, 'Cartesia Sonic-2  (TTS)  →  Audio', 9.5, True, '#8B0000')

# ═══════════════════════════════════════════════════════════════════════════════
# INTERNAL ARROWS (within agent.py)
# ═══════════════════════════════════════════════════════════════════════════════
arr(2.34, 2.29, 2.57, 2.29, '#556688', lw=1.6)   # VAD  → Whisper
arr(4.82, 2.29, 5.04, 2.29, '#556688', lw=1.6)   # Whisper → GPT
arr(7.17, 1.27, 7.17, 1.00, '#7B2DBE', lw=1.6)   # GPT → Cartesia

# ═══════════════════════════════════════════════════════════════════════════════
# EXTERNAL NUMBERED ARROWS
# ═══════════════════════════════════════════════════════════════════════════════

# ➊  Browser bottom-left → api.py top   (straight down)
arr(2.00, 5.75, 2.00, 5.45, lw=2.2)
badge(2.95, 5.60, '1. GET /getToken')

# ➋  api.py right → LiveKit left        (straight right)
arr(3.75, 4.50, 4.60, 4.50, lw=2.2)
badge(4.18, 4.72, '2. JWT token')

# ➌  LiveKit bottom → agent.py top (Whisper column)
arr(6.00, 3.35, 3.70, 2.76, lw=2.2, rad=0.18)
badge(4.60, 3.18, '3. Audio stream\n(user mic)')

# ➍  Cartesia right → LiveKit bottom-right (bot voice up)
arr(9.29, 0.64, 10.20, 3.35, lw=2.2, rad=-0.22)
badge(10.42, 2.00, '4. Audio stream\n(bot voice)')

# ➎  LiveKit top-right → Browser bottom-right
arr(10.20, 5.45, 12.15, 5.75, lw=2.2, rad=-0.18)
badge(11.42, 5.52, '5. Audio out')

# ── Save PNG ──────────────────────────────────────────────────────────────────
fig.savefig(OUT_PNG, dpi=150, bbox_inches='tight',
            facecolor='#F4F7FC', pad_inches=0.06)
plt.close(fig)
print(f'✓ PNG saved  →  {OUT_PNG}')

# ── Embed into PPTX ───────────────────────────────────────────────────────────
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.50)

sl = prs.slides.add_slide(prs.slide_layouts[6])   # blank layout
sl.background.fill.solid()
sl.background.fill.fore_color.rgb = RGBColor(0xF4, 0xF7, 0xFC)

# Slide title
txb = sl.shapes.add_textbox(Inches(0.3), Inches(0.05), Inches(12.73), Inches(0.58))
p   = txb.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = 'VoxCart — System Architecture'
run.font.size  = Pt(24)
run.font.bold  = True
run.font.color.rgb = RGBColor(0x1A, 0x2A, 0x6C)

# Architecture diagram image
sl.shapes.add_picture(OUT_PNG, Inches(0.0), Inches(0.63), Inches(13.33), Inches(6.87))

prs.save(OUT_PPTX)
print(f'✓ PPTX saved →  {OUT_PPTX}')
